"""
Material parser — extracts lecture text unit-by-unit from uploaded files.

  .pptx  -> one unit per SLIDE
  .pdf   -> one unit per PAGE
  .docx  -> one unit per SECTION (split on Heading paragraphs; falls back
            to fixed-size paragraph chunks if there are no headings)

Each parser returns: List[Tuple[unit_number, text]]

All parsers are optional: if a library is missing, a clear error is raised
so the caller can report it instead of crashing the server.
"""

import io
from typing import List, Tuple


def _require(module_name: str):
    raise RuntimeError(
        f"Required library '{module_name}' is not installed on the server. "
        f"Add it to requirements.txt and redeploy."
    )


def parse_pptx(data: bytes) -> List[Tuple[int, str]]:
    try:
        from pptx import Presentation
    except ImportError:
        _require("python-pptx")

    prs = Presentation(io.BytesIO(data))
    units: List[Tuple[int, str]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        parts.append(text)
            # Pull text from tables too
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
        # Speaker notes add useful context
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"(Notes) {notes}")
        except Exception:
            pass
        units.append((idx, "\n".join(parts).strip()))
    return units


def parse_pdf(data: bytes) -> List[Tuple[int, str]]:
    try:
        from pypdf import PdfReader
    except ImportError:
        _require("pypdf")

    reader = PdfReader(io.BytesIO(data))
    units: List[Tuple[int, str]] = []
    for idx, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        units.append((idx, text.strip()))
    return units


def parse_docx(data: bytes) -> List[Tuple[int, str]]:
    try:
        from docx import Document
    except ImportError:
        _require("python-docx")

    doc = Document(io.BytesIO(data))

    sections: List[List[str]] = []
    current: List[str] = []
    saw_heading = False

    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        style = (para.style.name or "").lower() if para.style else ""
        if style.startswith("heading"):
            saw_heading = True
            if current:
                sections.append(current)
            current = [text]
        else:
            current.append(text)
    if current:
        sections.append(current)

    # No headings? Fall back to chunks of ~8 paragraphs.
    if not saw_heading:
        all_paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        sections = [all_paras[i:i + 8] for i in range(0, len(all_paras), 8)] or [all_paras]

    return [(idx, "\n".join(sec).strip()) for idx, sec in enumerate(sections, start=1) if sec]


def parse_material(filename: str, data: bytes) -> Tuple[str, List[Tuple[int, str]]]:
    """
    Dispatch by file extension.
    Returns (unit_label, units) where unit_label is "slide" | "page" | "section".
    """
    name = (filename or "").lower()
    if name.endswith(".pptx"):
        return "slide", parse_pptx(data)
    if name.endswith(".pdf"):
        return "page", parse_pdf(data)
    if name.endswith(".docx"):
        return "section", parse_docx(data)
    raise ValueError(
        "Unsupported file type. Please upload a .pptx, .pdf, or .docx file."
    )
